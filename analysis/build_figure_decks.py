import json, re, os, html, datetime
from PIL import Image
ROOT="/nfsshare/users/P126156127/workspace/bioinf/gastric_cancer"
os.chdir(ROOT)
figs = json.load(open("/tmp/figs_flat.json"))

def key(f):
    m = re.match(r'Supplementary Figure S(\d+)([A-Z]?)', f['label'])
    if m: return (1, int(m.group(1)), m.group(2) or "")
    return (0, int(re.search(r'(\d+)', f['label']).group(1)), "")
figs.sort(key=key)

def split_cap(c):
    m = re.match(r'((?:Supplementary )?Figure S?\d+\.)\s*(.*)', c, re.S)
    return (m.group(1), m.group(2).strip()) if m else ("", c)

# ---------------- PDF (one image per page, no panelling) ----------------
CSS = """
@page { size:A4; margin:12mm; @bottom-center{content:counter(page);font-size:8pt;color:#888;} }
body { font-family:'DejaVu Sans',Helvetica,sans-serif; color:#111; }
.cover{text-align:center;padding-top:70mm;page-break-after:always;}
.cover h1{font-size:21pt;margin-bottom:3mm;} .cover .s{color:#555;font-size:10.5pt;}
.cover .n{margin-top:12mm;font-size:9pt;color:#777;}
.fig{page-break-after:always;} .fig:last-child{page-break-after:auto;}
.lab{font-size:11.5pt;font-weight:bold;margin:0 0 2mm 0;color:#0d3349;}
.img{text-align:center;margin-bottom:3mm;}
.img img{max-width:100%;max-height:210mm;object-fit:contain;}
.cap{font-size:8.3pt;line-height:1.45;text-align:justify;border-top:0.6px solid #ccc;padding-top:2mm;}
.src{font-size:7pt;color:#888;margin-top:1.2mm;font-family:'DejaVu Sans Mono',monospace;}
"""
today = datetime.date.today().strftime("%d %B %Y")
H=[f"""<div class='cover'><h1>Gastric Cancer Multi-Omics</h1>
<div class='s'>Complete figure set &mdash; every image standalone, one per page</div>
<div class='n'>{len(figs)} figures &middot; generated {today}<br>each reproduced from its committed source file</div></div>"""]
for f in figs:
    lab, rest = split_cap(f['caption'])
    abspath = os.path.join(ROOT, f['path'])
    H.append(f"""<div class='fig'><div class='lab'>{html.escape(f['label'])}</div>
<div class='img'><img src="file://{html.escape(abspath)}"></div>
<div class='cap'>{'<b>'+html.escape(lab)+'</b> ' if lab else ''}{html.escape(rest)}</div>
<div class='src'>{html.escape(f['path'])}</div></div>""")
open("/tmp/figs.html","w").write(f"<html><head><meta charset='utf-8'><style>{CSS}</style></head><body>{''.join(H)}</body></html>")

# ---------------- PPTX (editable, one image per slide) ----------------
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height; blank = prs.slide_layouts[6]

def tb(sl,l,t,w,h,txt,size,bold=False,color=(0x11,0x11,0x11),align=PP_ALIGN.LEFT):
    box = sl.shapes.add_textbox(l,t,w,h); tf = box.text_frame; tf.word_wrap=True
    p = tf.paragraphs[0]; p.alignment=align
    r = p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=RGBColor(*color); r.font.name="Calibri"

s = prs.slides.add_slide(blank)
tb(s,Inches(0.8),Inches(2.6),SW-Inches(1.6),Inches(1.0),"Gastric Cancer Multi-Omics",40,True,(0x0d,0x33,0x49),PP_ALIGN.CENTER)
tb(s,Inches(0.8),Inches(3.7),SW-Inches(1.6),Inches(0.6),f"Complete figure set — {len(figs)} standalone figures",18,False,(0x44,0x55,0x66),PP_ALIGN.CENTER)
tb(s,Inches(0.8),Inches(4.35),SW-Inches(1.6),Inches(0.5),f"Generated {today} · fully editable — move, resize or replace any element",12,False,(0x88,0x88,0x88),PP_ALIGN.CENTER)

MARG,TOPH,CAPH = Inches(0.45),Inches(0.60),Inches(1.30)
for f in figs:
    sl = prs.slides.add_slide(blank)
    lab, rest = split_cap(f['caption'])
    tb(sl,MARG,Inches(0.16),SW-2*MARG,TOPH,f['label'],21,True,(0x0d,0x33,0x49))
    aw, ah = SW-2*MARG, SH-TOPH-CAPH-Inches(0.30)
    iw, ih = Image.open(f['path']).size
    sc = min(aw/iw, ah/ih); w,h = int(iw*sc), int(ih*sc)
    sl.shapes.add_picture(f['path'], int((SW-w)/2), int(TOPH+Inches(0.16)), w, h)
    cap = (lab+" "+rest).strip()
    tb(sl,MARG,SH-CAPH-Inches(0.10),SW-2*MARG,CAPH,cap[:480],10.5,False,(0x22,0x22,0x22))
    tb(sl,MARG,SH-Inches(0.28),SW-2*MARG,Inches(0.24),f['path'],7.5,False,(0x99,0x99,0x99))
prs.save("FIGURES.pptx")
print("slides:", 1+len(figs), "| figures:", len(figs))
