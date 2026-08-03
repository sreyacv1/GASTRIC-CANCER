import json, os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime
os.chdir("/nfsshare/users/P126156127/workspace/bioinf/gastric_cancer")
figs = json.load(open("/tmp/figs_deck.json"))

prs = Presentation(); prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height; blank = prs.slide_layouts[6]
def tb(sl,l,t,w,h,txt,size,bold=False,color=(0x11,0x11,0x11),align=PP_ALIGN.LEFT):
    box=sl.shapes.add_textbox(l,t,w,h); tf=box.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    r=p.add_run(); r.text=txt; r.font.size=Pt(size); r.font.bold=bold
    r.font.color.rgb=RGBColor(*color); r.font.name="Calibri"

today=datetime.date.today().strftime("%d %B %Y")
s=prs.slides.add_slide(blank)
tb(s,Inches(0.8),Inches(2.5),SW-Inches(1.6),Inches(1.0),"Gastric Cancer Multi-Omics",40,True,(0x0d,0x33,0x49),PP_ALIGN.CENTER)
tb(s,Inches(0.8),Inches(3.6),SW-Inches(1.6),Inches(0.6),f"Complete figure archive — all {len(figs)} images, one per slide",18,False,(0x44,0x55,0x66),PP_ALIGN.CENTER)
tb(s,Inches(0.8),Inches(4.25),SW-Inches(1.6),Inches(0.5),f"Generated {today} · fully editable",12,False,(0x88,0x88,0x88),PP_ALIGN.CENTER)

MARG,TOPH,CAPH=Inches(0.45),Inches(0.58),Inches(1.05)
cur_section=None
for f in figs:
    if f["section"]!=cur_section:
        cur_section=f["section"]
        ds=prs.slides.add_slide(blank)
        tb(ds,Inches(0.8),Inches(3.1),SW-Inches(1.6),Inches(1.2),cur_section,32,True,(0x0d,0x33,0x49),PP_ALIGN.CENTER)
    sl=prs.slides.add_slide(blank)
    tag = f"{f['label']}" + (f"   ·   {f['label_hint']}" if f.get("label_hint") else "")
    tb(sl,MARG,Inches(0.16),SW-2*MARG,TOPH,tag,19,True,(0x0d,0x33,0x49))
    aw,ah = SW-2*MARG, SH-TOPH-CAPH-Inches(0.28)
    iw,ih = Image.open(f["path"]).size
    sc=min(aw/iw, ah/ih); w,h=int(iw*sc),int(ih*sc)
    sl.shapes.add_picture(f["path"], int((SW-w)/2), int(TOPH+Inches(0.16)), w, h)
    tb(sl,MARG,SH-CAPH-Inches(0.08),SW-2*MARG,CAPH,f["caption"][:430],10.5,False,(0x22,0x22,0x22))
    tb(sl,MARG,SH-Inches(0.27),SW-2*MARG,Inches(0.22),f["path"],7.5,False,(0x99,0x99,0x99))
prs.save("FIGURES_ALL.pptx")
n=len(prs.slides._sldIdLst)
print("slides:",n,"| images:",len(figs))
